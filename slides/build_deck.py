#!/usr/bin/env python3
"""
Builds SDR_Demo.pptx -- a Google-Slides-ready deck (upload to Drive, "Open with
Google Slides" and it converts to a native, editable Slides file).

Run:  .venv/bin/python3 build_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

ASSETS = Path(__file__).resolve().parent / "assets"

# ---------------------------------------------------------------- palette --
BG = RGBColor(0x0A, 0x19, 0x29)  # deep navy, "waterfall" background
BG_ALT = RGBColor(0x0F, 0x24, 0x38)  # slightly lighter navy for panels
INK = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0xD7, 0xE3, 0xEC)
MUTED = RGBColor(0x8E, 0xA7, 0xBA)
TEAL = RGBColor(0x35, 0xD0, 0xBA)   # accent -- waterfall highlight
TEAL_DIM = RGBColor(0x1B, 0x5E, 0x58)
AMBER = RGBColor(0xFF, 0xB0, 0x20)  # legal / caution accent
AMBER_DIM = RGBColor(0x4A, 0x36, 0x12)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ------------------------------------------------------------ primitives --
def new_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    # push background behind everything else added later
    s.shapes._spTree.remove(rect._element)
    s.shapes._spTree.insert(2, rect._element)
    return s


def add_text(slide, x, y, w, h, text, size=18, color=BODY, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT, line_spacing=1.15, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=18, color=BODY, bullet_color=TEAL,
                 gap=10, bold_lead=None):
    """items: list of str, or (str, sub_items) for one level of nesting."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        sub = None
        if isinstance(item, tuple):
            item, sub = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
        if sub:
            for s_item in sub:
                sp = tf.add_paragraph()
                sp.space_after = Pt(6)
                sp.level = 1
                sr = sp.add_run()
                sr.text = f"–  {s_item}"
                sr.font.size = Pt(size - 3)
                sr.font.color.rgb = MUTED
                sr.font.name = FONT
    return tb


def accent_bar(slide, y=Inches(1.5), color=TEAL, x=Inches(0.7), w=Inches(1.1)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def kicker_title(slide, kicker, title, kicker_color=TEAL, title_size=32):
    add_text(slide, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4),
              kicker.upper(), size=14, color=kicker_color, bold=True)
    add_text(slide, Inches(0.7), Inches(0.86), Inches(11.9), Inches(0.9),
              title, size=title_size, color=INK, bold=True)
    accent_bar(slide, y=Inches(1.62))


def page_number(slide, n, total):
    add_text(slide, Inches(12.35), Inches(7.08), Inches(0.8), Inches(0.3),
              f"{n} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def footer(slide, text="SDR Demo -- GNU Radio + HackRF + Ham Handheld"):
    add_text(slide, Inches(0.7), Inches(7.08), Inches(8), Inches(0.3),
              text, size=10, color=MUTED)


def pill(slide, x, y, w, h, text, fill=TEAL_DIM, text_color=TEAL, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.font.name = FONT
    return shp


TOTAL_SLIDES = 24
_n = [0]


def count():
    _n[0] += 1
    return _n[0]


# ============================================================== SLIDE 1 ==
s = new_slide()
add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.5),
          "SOFTWARE DEFINED RADIO", size=18, color=TEAL, bold=True)
add_text(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.6),
          "From Bits to Airwaves", size=48, color=INK, bold=True)
add_text(s, Inches(0.9), Inches(4.05), Inches(11.0), Inches(0.6),
          "A live demo with GNU Radio Companion, a HackRF One, and a 144/440 MHz handheld",
          size=20, color=BODY)
accent_bar(s, y=Inches(4.75), x=Inches(0.92), w=Inches(1.4))
add_text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4),
          "Mark Grennan  •  W5TSU  •  " + "2026",
          size=14, color=MUTED)

# ============================================================== SLIDE 2 ==
s = new_slide()
kicker_title(s, "Before We Touch The Radio", "The Airwaves: 3 Hz – 3 THz")
add_text(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.75),
          "By international convention (ITU), everything humans use for wireless "
          "communication -- from submarine navigation to satellite links -- lives "
          "somewhere in this range. Twelve decades, back to back, and every one of "
          "them is already spoken for.",
          size=18, color=BODY)
# band ladder: 12 ITU-designated decades, ELF (3 Hz) through THF (3 THz)
bands = ["ELF", "SLF", "ULF", "VLF", "LF", "MF", "HF", "VHF", "UHF", "SHF", "EHF", "THF"]
highlight = {"VHF", "UHF"}
ladder_y = Inches(3.05)
chip_w = Inches(0.90)
gap = Inches(0.10)
total_w = chip_w * len(bands) + gap * (len(bands) - 1)
lx = (SLIDE_W - total_w) / 2
for band in bands:
    is_hi = band in highlight
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ladder_y, chip_w, Inches(0.95))
    chip.adjustments[0] = 0.15
    chip.fill.solid()
    chip.fill.fore_color.rgb = TEAL if is_hi else BG_ALT
    chip.line.color.rgb = TEAL
    chip.line.width = Pt(1.5 if is_hi else 0.75)
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = band
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = BG if is_hi else INK
    r.font.name = FONT
    lx += chip_w + gap
add_text(s, Inches(0.7), Inches(4.15), Inches(2.5), Inches(0.35),
          "3 Hz", size=13, color=MUTED, bold=True)
add_text(s, Inches(10.13), Inches(4.15), Inches(2.5), Inches(0.35),
          "3 THz", size=13, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
add_text(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.35),
          "Each band spans one decade -- band N runs 0.3×10ᴺ Hz to 3×10ᴺ Hz.",
          size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
pill(s, Inches(3.9), Inches(5.35), Inches(5.5), Inches(0.45),
     "Tonight lives here: VHF / UHF", fill=TEAL_DIM, text_color=TEAL, size=15)
add_bullets(s, Inches(1.6), Inches(6.0), Inches(10.1), Inches(1.0), [
    "2m (144-148 MHz) and 70cm (420-450 MHz) -- our ham bands -- plus FM "
    "broadcast (88-108 MHz), all three narrow slivers of one 270:1 span (30 "
    "MHz-3 GHz).",
], size=16, gap=0)
add_text(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
          "ITU Radio Regulations, Article 2 -- band designations ELF-THF",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 3 ==
s = new_slide()
kicker_title(s, "Before We Touch The Radio",
             "Every Signal Arrives Already Mixed With Noise")
p1x = Inches(0.7); pw = Inches(5.7)
panel1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p1x, Inches(2.0), pw, Inches(4.15))
panel1.adjustments[0] = 0.04
panel1.fill.solid(); panel1.fill.fore_color.rgb = BG_ALT
panel1.line.color.rgb = TEAL_DIM; panel1.line.width = Pt(1)
panel1.shadow.inherit = False
add_text(s, p1x + Inches(0.35), Inches(2.25), pw - Inches(0.7), Inches(0.4),
          "Noise From Nature", size=19, color=TEAL, bold=True)
add_bullets(s, p1x + Inches(0.35), Inches(2.75), pw - Inches(0.7), Inches(3.3), [
    "Thermal noise: the atoms in your receiver's own components are "
    "vibrating with heat -- that's noise, before an antenna is even "
    "involved.",
    "Modeled as Gaussian / \"white\" noise -- flat across the whole band. "
    "That's the hiss you hear with nothing tuned in.",
    "Lightning sferics: crackly pops from distant thunderstorms, audible "
    "from VLF up through HF.",
    "Ionized meteor trails briefly reflect distant signals into a "
    "receiver -- noise and signal from the same natural event.",
], size=15, gap=10)

p2x = Inches(6.9)
panel2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p2x, Inches(2.0), pw, Inches(4.15))
panel2.adjustments[0] = 0.04
panel2.fill.solid(); panel2.fill.fore_color.rgb = BG_ALT
panel2.line.color.rgb = TEAL_DIM; panel2.line.width = Pt(1)
panel2.shadow.inherit = False
add_text(s, p2x + Inches(0.35), Inches(2.25), pw - Inches(0.7), Inches(0.4),
          "Noise We Make Ourselves", size=19, color=TEAL, bold=True)
add_bullets(s, p2x + Inches(0.35), Inches(2.75), pw - Inches(0.7), Inches(3.3), [
    "Every switching power supply, phone charger, and LED driver in this "
    "room is radiating RF -- one of the most common interference sources "
    "hams fight.",
    "USB 3.0 leaks clock noise at multiples of 480 MHz; HDMI cables can "
    "splatter interference across roughly 148-742 MHz.",
    "Monitors, DisplayPort, powerline networking adapters -- all of it "
    "adds more.",
    "None of it is malicious. It's just what happens when fast digital "
    "clocks run next to a receiver.",
], size=15, gap=10)
add_text(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.65),
          "Every demo tonight starts here: find the signal you want inside "
          "everything else that's already there. Watch the noise floor on the "
          "very first demo -- it's not empty.",
          size=16, color=TEAL, italic=True)
add_text(s, Inches(0.7), Inches(7.08), Inches(11), Inches(0.3),
          "Sources: pysdr.org (Ch. 10, Noise and Random Variables) and "
          "sigidwiki.com (Interfering Emissions category)",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 4 ==
s = new_slide()
kicker_title(s, "Before We Touch The Radio", "Hear It For Yourself, Live")
pill(s, Inches(0.7), Inches(1.85), Inches(1.5), Inches(0.35), "RECEIVE ONLY", fill=TEAL_DIM, text_color=TEAL, size=12)
add_bullets(s, Inches(0.7), Inches(2.45), Inches(11.9), Inches(1.6), [
    "Everything on the last slide -- thermal noise, atmospheric sferics, "
    "switching power supplies -- is audible right now, live.",
    "The entire 20-meter ham band (14.000-14.350 MHz), captured all at "
    "once and played as one audio stream -- not tuned to any single "
    "station, just the raw mix of everything active on the band.",
], size=17, gap=12)
add_text(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(0.35),
          "Prefer to go exploring on your own later? A worldwide network of "
          "live receivers anyone can tune:",
          size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
link_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.67), Inches(4.8), Inches(6.0), Inches(0.75))
link_btn.adjustments[0] = 0.25
link_btn.fill.solid(); link_btn.fill.fore_color.rgb = TEAL_DIM
link_btn.line.color.rgb = TEAL; link_btn.line.width = Pt(1.5)
link_btn.shadow.inherit = False
link_btn_tf = link_btn.text_frame
link_btn_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
link_btn_p = link_btn_tf.paragraphs[0]
link_btn_p.alignment = PP_ALIGN.CENTER
link_btn_r = link_btn_p.add_run()
link_btn_r.text = "websdr.org  --  listen live"
link_btn_r.font.size = Pt(19); link_btn_r.font.bold = True
link_btn_r.font.color.rgb = TEAL; link_btn_r.font.name = FONT
link_btn_r.hyperlink.address = "http://websdr.org/"
add_text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.4),
          "Try it: pick any frequency and just listen.",
          size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(7.08), Inches(11), Inches(0.3),
          "websdr.org -- University of Twente WebSDR, public since April 2008 "
          "(Pieter-Tjerk de Boer, PA3FWM)",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 5 ==
s = new_slide()
kicker_title(s, "Before We Touch The Radio", "Finding a Signal in the Noise")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(1.9), [
    "A receiver doesn't remove noise -- it narrows down. Filter to just "
    "the channel bandwidth of the signal you want, and every watt of "
    "noise outside that channel disappears with it.",
    "What actually matters is SNR: not \"is there noise\" (there always "
    "is) but \"is the signal stronger than the noise inside the "
    "bandwidth you're listening in.\"",
    "Modulation is what makes a signal findable in the first place -- "
    "FM, AM, PSK, and friends each carve a distinct, structured shape "
    "into the spectrum that a demodulator can lock onto and reject "
    "everything else.",
], size=16, gap=12)
# mini diagram: Wideband capture (mostly noise) -> Filter to channel -> Demodulate
fs_dy = Inches(4.35)
fs_items = ["Wideband Capture\n(mostly noise)", "Filter to\nChannel Bandwidth", "Demodulate\n(now readable)"]
fs_bw = Inches(3.3); fs_gap = Inches(0.55)
fs_total_w = fs_bw * len(fs_items) + fs_gap * (len(fs_items) - 1)
fs_bx = (SLIDE_W - fs_total_w) / 2
for i, item in enumerate(fs_items):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fs_bx, fs_dy, fs_bw, Inches(1.1))
    chip.adjustments[0] = 0.12
    chip.fill.solid()
    chip.fill.fore_color.rgb = TEAL_DIM if i == 2 else BG_ALT
    chip.line.color.rgb = TEAL
    chip.line.width = Pt(1.25)
    chip.shadow.inherit = False
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = item
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    if i < len(fs_items) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, fs_bx + fs_bw, fs_dy + Inches(0.38), fs_gap, Inches(0.34))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background(); arrow.shadow.inherit = False
    fs_bx += fs_bw + fs_gap
add_text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.35),
          "Once you can pick a signal out of the noise, the next question is "
          "\"what am I even listening to\" -- a crowdsourced field guide:",
          size=14, color=BODY, align=PP_ALIGN.CENTER)
link_tb5 = s.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.4))
link_tf5 = link_tb5.text_frame
link_p5 = link_tf5.paragraphs[0]
link_p5.alignment = PP_ALIGN.CENTER
link_r5 = link_p5.add_run()
link_r5.text = "sigidwiki.com -- Signal Identification Guide"
link_r5.font.size = Pt(18); link_r5.font.bold = True
link_r5.font.color.rgb = TEAL; link_r5.font.name = FONT
link_r5.hyperlink.address = "https://www.sigidwiki.com/wiki/Signal_Identification_Guide"
add_text(s, Inches(0.7), Inches(7.08), Inches(11), Inches(0.3),
          "sigidwiki.com -- ~600 identified signals with example sounds and "
          "waterfall images, from VLF to UHF",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 6 ==
s = new_slide()
kicker_title(s, "The Idea", "What Is Software Defined Radio?")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6.9), Inches(4.5), [
    "Traditional radios: filtering, mixing, and demodulation are done by "
    "dedicated analog/digital hardware built for one job.",
    "SDR: an antenna feeds a wideband ADC/DAC, and everything after that -- "
    "tuning, filtering, demodulating, decoding -- happens in software.",
    "Same box, endless radios: change the software, not the hardware.",
    "Today's demo does this literally -- one HackRF, six different "
    "“radios,” just by loading a different flowgraph.",
], size=19, gap=16)
# simple diagram: Antenna -> ADC/DAC -> Software
panel_y = Inches(2.1)
dx = Inches(7.9)
dw = Inches(4.7)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dx, panel_y, dw, Inches(4.3))
box.adjustments[0] = 0.04
box.fill.solid(); box.fill.fore_color.rgb = BG_ALT
box.line.color.rgb = TEAL_DIM; box.line.width = Pt(1)
box.shadow.inherit = False
labels = ["Antenna", "ADC / DAC\n(HackRF)", "Software\n(GNU Radio)"]
sub = ["Analog RF in/out", "Samples ↔ RF", "Filter · demod · decode"]
by = panel_y + Inches(0.35)
for lab, sb in zip(labels, sub):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, dx + Inches(0.5), by, dw - Inches(1.0), Inches(0.95))
    chip.adjustments[0] = 0.12
    chip.fill.solid(); chip.fill.fore_color.rgb = TEAL_DIM
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lab
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sb
    r2.font.size = Pt(11); r2.font.color.rgb = MUTED; r2.font.name = FONT
    if lab != labels[-1]:
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, dx + dw/2 - Inches(0.15), by + Inches(0.95), Inches(0.3), Inches(0.25))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background()
        arrow.shadow.inherit = False
    by += Inches(1.2)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 7 ==
s = new_slide()
kicker_title(s, "The Hardware", "HackRF One")
# two panel layout: core specs vs. TX capabilities
p1x = Inches(0.7); pw = Inches(5.7)
panel1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p1x, Inches(2.1), pw, Inches(4.5))
panel1.adjustments[0] = 0.04
panel1.fill.solid(); panel1.fill.fore_color.rgb = BG_ALT
panel1.line.color.rgb = TEAL_DIM; panel1.line.width = Pt(1)
panel1.shadow.inherit = False
add_text(s, p1x + Inches(0.35), Inches(2.35), pw - Inches(0.7), Inches(0.4),
          "Core Specs", size=20, color=TEAL, bold=True)
add_bullets(s, p1x + Inches(0.35), Inches(2.85), pw - Inches(0.7), Inches(3.5), [
    "1 MHz – 6 GHz tuning range",
    "Up to 20 Msps, 8-bit quadrature ADC/DAC (8-bit I, 8-bit Q)",
    "Half-duplex (RX or TX, never both at once)",
    "Hi-Speed USB 2.0, USB-powered, open source hardware",
    "This talk: driven via SoapySDR from GNU Radio",
], size=16, gap=12)

p2x = Inches(6.9)
panel2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p2x, Inches(2.1), pw, Inches(4.5))
panel2.adjustments[0] = 0.04
panel2.fill.solid(); panel2.fill.fore_color.rgb = BG_ALT
panel2.line.color.rgb = TEAL_DIM; panel2.line.width = Pt(1)
panel2.shadow.inherit = False
add_text(s, p2x + Inches(0.35), Inches(2.35), pw - Inches(0.7), Inches(0.4),
          "TX Capabilities", size=20, color=TEAL, bold=True)
add_bullets(s, p2x + Inches(0.35), Inches(2.85), pw - Inches(0.7), Inches(3.5), [
    "Two software-controlled TX gain stages -- no physical knobs: an "
    "on/off RF amp (~+11 dB) and a VGA/IF gain stage, 0-47 dB in 1 dB "
    "steps.",
    "Max output power varies with frequency -- roughly +5 to +15 dBm "
    "from 1 MHz-2.17 GHz, dropping to about 0-10 dBm up through 4 GHz. "
    "Best performance is actually 2170-2740 MHz.",
    "Tonight's TX demos default to just 10 dB of VGA gain -- comfortably "
    "at the low end, per the \"minimum power necessary\" guidance.",
    "No built-in per-band TX filtering -- it's a general-purpose "
    "transmitter, so staying inside our allocated ham segments is on us.",
], size=15, gap=10)
add_text(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
          "Sources: hackrf.readthedocs.io, greatscottgadgets.com/hackrf/one",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 8 ==
s = new_slide()
kicker_title(s, "The Software", "GNU Radio + GNU Radio Companion")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6.7), Inches(4.5), [
    "GNU Radio: open-source DSP toolkit -- filters, modulators, "
    "demodulators, resamplers, all as reusable building blocks.",
    "GNU Radio Companion (GRC): drag-and-drop flowgraph editor that "
    "generates the underlying Python for you.",
    "SoapySDR: hardware abstraction layer -- the same flowgraph blocks "
    "work whether it's a HackRF, RTL-SDR, LimeSDR, or others.",
    "Every flowgraph in this demo is a plain-text .grc file -- versionable, "
    "diffable, shareable.",
], size=18, gap=16)
img_x = Inches(7.9); img_w = Inches(4.7)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_x, Inches(2.1), img_w, Inches(4.3))
panel.adjustments[0] = 0.04
panel.fill.solid(); panel.fill.fore_color.rgb = BG_ALT
panel.line.color.rgb = TEAL_DIM; panel.line.width = Pt(1)
panel.shadow.inherit = False
add_text(s, img_x + Inches(0.35), Inches(2.35), img_w - Inches(0.7), Inches(0.35),
          "Toolchain", size=15, color=MUTED, bold=True)
chain = ["Flowgraph (.grc)", "GNU Radio Companion", "Generated Python", "GNU Radio runtime + SoapySDR", "HackRF One"]
cy = Inches(2.85)
for i, item in enumerate(chain):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_x + Inches(0.35), cy, img_w - Inches(0.7), Inches(0.62))
    chip.adjustments[0] = 0.2
    chip.fill.solid(); chip.fill.fore_color.rgb = TEAL_DIM
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = item
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    if i < len(chain) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, img_x + img_w/2 - Inches(0.12), cy + Inches(0.62), Inches(0.24), Inches(0.16))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background(); arrow.shadow.inherit = False
    cy += Inches(0.78)
footer(s); page_number(s, count(), TOTAL_SLIDES)
s.shapes.add_picture(str(ASSETS / "gnuradio_logo.png"),
                      Inches(8.51), Inches(0.24), Inches(3.66), Inches(1.83))

# ============================================================== SLIDE 9 ==
s = new_slide()
kicker_title(s, "Core Concept", "Reading a Flowgraph")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.7), Inches(2.0), [
    "Blocks are connected left to right: a Source produces samples, "
    "middle blocks process them, a Sink consumes them.",
    "Everything on the RF side flows as complex IQ samples -- one number "
    "for signal strength/phase in two dimensions (I and Q), captured at "
    "the sample rate.",
    "Sample rate sets your instantaneous bandwidth: capture at 6 Msps and "
    "you're watching 6 MHz of spectrum at once, not just one channel.",
], size=18, gap=14)
# mini diagram: Source -> Filter -> Demod -> Sink
dy = Inches(4.6)
items = ["HackRF\nSource", "Filter /\nResample", "Demodulate", "Audio / File\nSink"]
bw = Inches(2.55); gap_w = Inches(0.55)
total_w = bw * len(items) + gap_w * (len(items) - 1)
start_x = (SLIDE_W - total_w) / 2
bx = start_x
for i, item in enumerate(items):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, dy, bw, Inches(1.1))
    chip.adjustments[0] = 0.12
    chip.fill.solid()
    chip.fill.fore_color.rgb = TEAL_DIM if i in (0, 3) else BG_ALT
    chip.line.color.rgb = TEAL
    chip.line.width = Pt(1.25)
    chip.shadow.inherit = False
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = item
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    if i < len(items) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, bx + bw, dy + Inches(0.38), gap_w, Inches(0.34))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background(); arrow.shadow.inherit = False
    bx += bw + gap_w
footer(s); page_number(s, count(), TOTAL_SLIDES)
s.shapes.add_picture(str(ASSETS / "grc_flowgraph_screenshot.jpg"),
                      Inches(6.94), Inches(0.19), Inches(5.20), Inches(1.64))

# ============================================================== SLIDE 10 ==
s = new_slide()
kicker_title(s, "Demo 0  ·  Icebreaker", "FM Broadcast Receiver", kicker_color=MUTED)
pill(s, Inches(0.7), Inches(1.85), Inches(1.5), Inches(0.35), "RECEIVE ONLY", fill=TEAL_DIM, text_color=TEAL, size=12)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(11.7), Inches(3.6), [
    "A quick tune to any strong local FM station brings music through the "
    "laptop speakers, instantly.",
    "Not a ham-band demo -- it proves the whole chain works (driver, "
    "HackRF, GNU Radio, audio out) before the ham material, and it needs "
    "no license to receive.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================== SLIDE 11 ==
s = new_slide()
kicker_title(s, "Demo 1", "Live Spectrum & Waterfall", kicker_color=TEAL)
pill(s, Inches(0.7), Inches(1.85), Inches(1.5), Inches(0.35), "RECEIVE ONLY", fill=TEAL_DIM, text_color=TEAL, size=12)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(11.7), Inches(3.8), [
    "One flowgraph covers both ham bands -- we'll retune live between 2m "
    "and 70cm.",
    "Watch the spectrum plot and waterfall light up in real time as a "
    "handheld transmission appears.",
    "It's the same slice of RF spectrum a traditional radio can only "
    "look at one narrow channel of at a time.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 12 ==
s = new_slide()
kicker_title(s, "Demo 2", "Tunable NBFM Voice Receiver", kicker_color=TEAL)
pill(s, Inches(0.7), Inches(1.85), Inches(1.5), Inches(0.35), "RECEIVE ONLY", fill=TEAL_DIM, text_color=TEAL, size=12)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(11.7), Inches(3.8), [
    "A real narrowband FM voice receiver -- channel filter, resampler, "
    "FM demodulator, straight to the speakers.",
    "Tuned to 146.520 MHz, the national 2m simplex calling frequency.",
    "This whole receiver -- normally a dedicated chip in an off-the-shelf "
    "radio -- is about a dozen blocks of software here.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 13 ==
s = new_slide(bg=RGBColor(0x1A, 0x14, 0x06))
add_text(s, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.4),
          "⚠  BEFORE WE TRANSMIT", size=16, color=AMBER, bold=True)
add_text(s, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.9),
          "Part 97 Ground Rules", size=34, color=INK, bold=True)
accent_bar(s, y=Inches(1.85), color=AMBER)
add_bullets(s, Inches(0.7), Inches(2.25), Inches(11.7), Inches(4.5), [
    "A licensed control operator is present the entire time we transmit.",
    "We identify with callsign at the start/end and every 10 minutes.",
    "No music, no broadcast content -- a brief, identified test "
    "transmission only.",
    "Power stays at the minimum needed -- TX gain starts low and only "
    "comes up as far as the demo needs.",
    "No control operator today, or want zero on-air emission? The HackRF "
    "TX port can go into a dummy load instead of an antenna -- same demo, "
    "nothing radiated.",
], size=18, color=RGBColor(0xF3, 0xE3, 0xC7), gap=14)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 14 ==
s = new_slide()
add_text(s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4),
          "DEMO 3 · THE HEADLINE DEMO".upper(), size=14, color=AMBER, bold=True)
add_text(s, Inches(0.7), Inches(0.86), Inches(11.9), Inches(1.0),
          "Record the Entire 2m Band", size=34, color=INK, bold=True)
accent_bar(s, y=Inches(1.7), color=AMBER)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.7), Inches(4.5), [
    "Tunes to 146.0 MHz center at 6 Msps -- covers 143.0 to 149.0 MHz, the "
    "whole 144-148 MHz 2m band with margin to spare.",
    ("Not one channel -- every signal on the entire band, captured at once",
     ["Raw IQ, written to a self-describing file (sample rate + center "
      "frequency saved in the header)"]),
    "A single recording captures everything transmitted on the band "
    "during that window -- roughly 2.9 GB per minute at this sample rate.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 15 ==
s = new_slide()
add_text(s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.4),
          "DEMO 3B · THE PAYOFF".upper(), size=14, color=AMBER, bold=True)
add_text(s, Inches(0.7), Inches(0.86), Inches(11.9), Inches(1.0),
          "...And Play It Back Over the Air", size=34, color=INK, bold=True)
accent_bar(s, y=Inches(1.7), color=AMBER)
pill(s, Inches(0.7), Inches(1.9), Inches(2.3), Inches(0.4), "TRANSMITS — SEE NEXT SLIDE",
     fill=AMBER_DIM, text_color=AMBER, size=12)
add_bullets(s, Inches(0.7), Inches(2.55), Inches(11.7), Inches(4.0), [
    "Reads the file back and feeds it straight into the HackRF's "
    "transmitter -- the exact recorded band reappears on the air.",
    "The waterfall reproduces exactly what was captured; the handheld "
    "hears the same signal a second time, live.",
    "This is the moment that makes SDR click: the recording *is* the "
    "radio signal -- numbers on disk, transmitted back as RF.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 16 ==
s = new_slide()
kicker_title(s, "Demo 4  ·  Optional", "NBFM Voice Transmitter", kicker_color=AMBER)
pill(s, Inches(0.7), Inches(1.85), Inches(2.5), Inches(0.35), "TRANSMITS — CONTROL OP REQUIRED", fill=AMBER_DIM, text_color=AMBER, size=12)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(11.7), Inches(3.8), [
    "The reverse direction: laptop microphone → NBFM modulator → HackRF "
    "→ antenna.",
    "Live two-way on the same 146.520 MHz simplex frequency as Demo 2.",
], size=19, gap=16)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 17 ==
s = new_slide()
kicker_title(s, "The Software, Revisited", "GNU Radio: What It Is, How It's Used")
p1x = Inches(0.7); pw = Inches(5.7)
panel1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p1x, Inches(2.1), pw, Inches(4.5))
panel1.adjustments[0] = 0.04
panel1.fill.solid(); panel1.fill.fore_color.rgb = BG_ALT
panel1.line.color.rgb = TEAL_DIM; panel1.line.width = Pt(1)
panel1.shadow.inherit = False
add_text(s, p1x + Inches(0.35), Inches(2.35), pw - Inches(0.7), Inches(0.4),
          "What It Actually Is", size=19, color=TEAL, bold=True)
add_bullets(s, p1x + Inches(0.35), Inches(2.85), pw - Inches(0.7), Inches(3.5), [
    "Started in 2001 (Eric Blossom); now stewarded by the GNU Radio "
    "Project and the GNU Radio Foundation.",
    "C++ underneath for real-time DSP performance; Python (and GNU Radio "
    "Companion) on top for building and scripting flowgraphs.",
    "A scheduler manages sample buffers between blocks -- the same "
    "runtime whether the samples are live from a HackRF or read back "
    "from a file.",
    "Hundreds of built-in blocks -- filters, modulators, channel coding, "
    "synchronization -- plus a large ecosystem of out-of-tree (OOT) "
    "modules for specific protocols.",
], size=15, gap=10)

p2x = Inches(6.9)
add_text(s, p2x + Inches(0.35), Inches(2.35), Inches(5.0), Inches(0.32),
          "How to learn GNU-Radio Companion", size=19, color=TEAL, bold=True)
link_tb = s.shapes.add_textbox(p2x + Inches(0.35), Inches(2.85), Inches(5.0), Inches(1.06))
link_tf = link_tb.text_frame
link_tf.word_wrap = True
link_p = link_tf.paragraphs[0]
link_r = link_p.add_run()
link_r.text = "https://www.youtube.com/playlist?list=PLywxmTaHNUNyKmgF70q8q3QHYIw_LFbrX"
link_r.font.size = Pt(15)
link_r.font.name = FONT
link_r.hyperlink.address = "https://www.youtube.com/playlist?list=PLywxmTaHNUNyKmgF70q8q3QHYIw_LFbrX"
add_text(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
          "Source: gnuradio.org -- GNU Radio Conference (GRCon) happens every year",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)
s.shapes.add_picture(str(ASSETS / "learn_gnuradio_youtube_thumb.png"),
                      Inches(6.93), Inches(3.61), Inches(5.76), Inches(3.43))

# ============================================================= SLIDE 18 ==
s = new_slide()
kicker_title(s, "Beyond Tonight's Setup", "Connecting GNU Radio to Other Radios")
p1x = Inches(0.7); pw = Inches(5.7)
panel1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p1x, Inches(2.1), pw, Inches(4.5))
panel1.adjustments[0] = 0.04
panel1.fill.solid(); panel1.fill.fore_color.rgb = BG_ALT
panel1.line.color.rgb = TEAL_DIM; panel1.line.width = Pt(1)
panel1.shadow.inherit = False
add_text(s, p1x + Inches(0.35), Inches(2.35), pw - Inches(0.7), Inches(0.4),
          "Other SDRs", size=19, color=TEAL, bold=True)
add_bullets(s, p1x + Inches(0.35), Inches(2.85), pw - Inches(0.7), Inches(3.5), [
    "Same GRC blocks, different hardware -- SoapySDR supports RTL-SDR, "
    "LimeSDR, PlutoSDR, Airspy, SDRplay, and more through one common "
    "driver API.",
    "USRP hardware typically goes through UHD instead -- GNU Radio ships "
    "dedicated UHD blocks for it.",
    "Swapping hardware is usually a one-block change in the flowgraph, "
    "not a redesign.",
], size=16, gap=12)

p2x = Inches(6.9)
panel2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, p2x, Inches(2.1), pw, Inches(4.5))
panel2.adjustments[0] = 0.04
panel2.fill.solid(); panel2.fill.fore_color.rgb = BG_ALT
panel2.line.color.rgb = TEAL_DIM; panel2.line.width = Pt(1)
panel2.shadow.inherit = False
add_text(s, p2x + Inches(0.35), Inches(2.35), pw - Inches(0.7), Inches(0.4),
          "Analog / Non-SDR Radios", size=19, color=TEAL, bold=True)
add_bullets(s, p2x + Inches(0.35), Inches(2.85), pw - Inches(0.7), Inches(3.5), [
    "Audio interface: a USB sound card wired to a radio's mic/speaker "
    "(or data) port runs digital modes -- APRS, PSK31, and more -- "
    "through a normal transceiver, no RF-direct hardware needed.",
    "Rig control (CAT): Hamlib-based blocks let a flowgraph key PTT and "
    "set frequency on a physical radio over serial/USB.",
    "Direct RF coupling: exactly what tonight's demo did -- the SDR "
    "transmits or receives over the air, and the other radio just needs "
    "an antenna.",
], size=16, gap=12)
add_text(s, Inches(0.7), Inches(7.08), Inches(11), Inches(0.3),
          "Tonight only used one SDR and one handheld over RF -- the same tools "
          "bridge to almost anything with an antenna or an audio jack.",
          size=13, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 19 ==
s = new_slide()
kicker_title(s, "One Word, Three Meanings", "Three Kinds of Bandwidth")
bw_panels = [
    ("SDR (RF) Bandwidth", [
        "Instantaneous bandwidth = sample rate itself -- spans -Fs/2 to "
        "+Fs/2 around wherever you're tuned.",
        "HackRF: up to 20 Msps -> up to 20 MHz of spectrum visible at "
        "once.",
        "Rule of thumb: trust only the center ~4/5 of that (\"Sean's 4/5 "
        "rule\") -- the anti-alias filter rolls off near the edges.",
    ]),
    ("Interface Bandwidth (USB)", [
        "A separate bottleneck: getting every sample off the SDR and "
        "into the computer, in real time, without dropping any.",
        "HackRF: USB 2.0 Hi-Speed. At 20 Msps complex 8-bit, that's "
        "~320 Mbps -- already close to USB2's practical ceiling.",
        "Exactly why hackrf_info warns about other devices sharing the "
        "USB bus at high sample rates.",
    ]),
    ("Receive Bandwidth (What the Software Decodes)", [
        "The RF capture and the channel actually being decoded are two "
        "different numbers.",
        "Demo 3 captured 6 MHz of the 2m band; one NBFM voice channel "
        "is only ~16 kHz wide.",
        "That gap is exactly the Filter/Resample block from \"Reading a "
        "Flowgraph\" -- narrowing a wide capture down to one channel.",
    ]),
]
bw_pw = Inches(3.83); bw_gap = Inches(0.2)
bw_x = Inches(0.7)
for heading, bullets in bw_panels:
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bw_x, Inches(2.1), bw_pw, Inches(4.5))
    panel.adjustments[0] = 0.05
    panel.fill.solid(); panel.fill.fore_color.rgb = BG_ALT
    panel.line.color.rgb = TEAL_DIM; panel.line.width = Pt(1)
    panel.shadow.inherit = False
    add_text(s, bw_x + Inches(0.25), Inches(2.3), bw_pw - Inches(0.5), Inches(0.7),
              heading, size=16, color=TEAL, bold=True)
    add_bullets(s, bw_x + Inches(0.25), Inches(3.05), bw_pw - Inches(0.5), Inches(3.4),
                bullets, size=13, gap=10)
    bw_x += bw_pw + bw_gap
add_text(s, Inches(0.7), Inches(7.08), Inches(11), Inches(0.3),
          "Source: pysdr.org (Ch. 3, IQ Sampling -- sample rate = instantaneous "
          "bandwidth, Nyquist, \"Sean's 4/5 rule\")",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 20 ==
s = new_slide()
kicker_title(s, "Where This Goes Next", "More Demo Ideas")
left_items = [
    "CTCSS/PL tone squelch decode",
    "CW (Morse) ID transmitter",
    "RDS decode from FM broadcast",
    "APRS decode (144.390 MHz)",
]
right_items = [
    "WSPR / FT8 weak-signal decode",
    "ADS-B aircraft tracking (1090 MHz)",
    "Cross-band digital repeater",
    "Digital voice (DMR / D-STAR) decode",
]
add_text(s, Inches(0.7), Inches(2.05), Inches(5.6), Inches(0.35), "Quick wins", size=16, color=TEAL, bold=True)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(5.6), Inches(3.8), left_items, size=17, gap=14)
add_text(s, Inches(6.9), Inches(2.05), Inches(5.6), Inches(0.35), "Bigger swings", size=16, color=TEAL, bold=True)
add_bullets(s, Inches(6.9), Inches(2.5), Inches(5.6), Inches(3.8), right_items, size=17, gap=14)
add_text(s, Inches(0.7), Inches(6.45), Inches(11), Inches(0.3),
          "Full writeup with effort/dependency notes: BRAINSTORM.md", size=13, color=MUTED, italic=True)
link_tb18 = s.shapes.add_textbox(Inches(0.7), Inches(6.78), Inches(11.9), Inches(0.3))
link_tf18 = link_tb18.text_frame
link_tf18.word_wrap = True
link_p18 = link_tf18.paragraphs[0]
r18a = link_p18.add_run()
r18a.text = "Working GNU Radio flowgraphs for many of these modes (ham club tutorial repo): "
r18a.font.size = Pt(13); r18a.font.italic = True; r18a.font.color.rgb = MUTED; r18a.font.name = FONT
r18b = link_p18.add_run()
r18b.text = "github.com/argilo/sdr-examples"
r18b.font.size = Pt(13); r18b.font.italic = True; r18b.font.name = FONT
r18b.hyperlink.address = "https://github.com/argilo/sdr-examples"
page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 21 ==
s = new_slide()
kicker_title(s, "Reference", "Specs at a Glance")
rows = [
    ("Frequency range", "1 MHz – 6 GHz (HackRF One)"),
    ("Sample rate used today", "2–10 Msps depending on demo"),
    ("2m band", "144.000 – 148.000 MHz"),
    ("70cm band", "420.000 – 450.000 MHz"),
    ("2m simplex calling freq", "146.520 MHz"),
    ("Demo 3 capture", "146.0 MHz center @ 6 Msps ≈ 2.9 GB/min"),
    ("HackRF TX power (typical)", "≈ 10–15 dBm max, kept low for these demos"),
]
tx = Inches(0.7); ty = Inches(2.1); tw = Inches(11.9); rh = Inches(0.62)
for i, (k, v) in enumerate(rows):
    row_bg = BG_ALT if i % 2 == 0 else BG
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tx, ty + rh * i, tw, rh)
    rect.fill.solid(); rect.fill.fore_color.rgb = row_bg
    rect.line.fill.background(); rect.shadow.inherit = False
    add_text(s, tx + Inches(0.25), ty + rh * i, Inches(4.5), rh, k, size=15,
              color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, tx + Inches(5.0), ty + rh * i, tw - Inches(5.2), rh, v, size=15,
              color=BODY, anchor=MSO_ANCHOR.MIDDLE)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 22 ==
s = new_slide()
kicker_title(s, "Resources", "Go Build Something")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.7), Inches(4.5), [
    ("GNU Radio", ["gnuradio.org — the toolkit, docs, and tutorials"]),
    ("HackRF One", ["greatscottgadgets.com/hackrf — hardware + firmware"]),
    ("SoapySDR", ["github.com/pothosware/SoapySDR — the hardware abstraction layer used today"]),
    ("ARRL", ["arrl.org — licensing info and Part 97 reference"]),
    ("This kit", ["All flowgraphs, README, and brainstorm doc are in this repo’s "
                   "flowgraphs/ and top-level README.md / BRAINSTORM.md"]),
], size=19, gap=18)
footer(s); page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 23 ==
s = new_slide()
kicker_title(s, "Software Defined Radio Requires Software", "Learn the DSP, in Python: pysdr.org")
add_bullets(s, Inches(0.7), Inches(2.1), Inches(11.7), Inches(1.15), [
    "Everything tonight ran through GNU Radio Companion's drag-and-drop "
    "canvas -- pysdr.org teaches the same DSP by writing the Python "
    "(NumPy/SciPy) yourself, concept by concept.",
    "Free, written for engineers, no textbook required first -- the same "
    "site cited earlier tonight for noise (slide 3) and bandwidth/Nyquist "
    "(slide 17).",
], size=16, gap=10)
# horizontal roadmap: FFT -> Filters -> Digital Modulation -> RX/TX in Python
py_dy = Inches(3.7)
py_items = ["Fourier\nTransforms", "Filters", "Digital\nModulation", "RX & TX\nin Python"]
py_bw = Inches(2.55); py_gap = Inches(0.55)
py_total_w = py_bw * len(py_items) + py_gap * (len(py_items) - 1)
py_bx = (SLIDE_W - py_total_w) / 2
for i, item in enumerate(py_items):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, py_bx, py_dy, py_bw, Inches(1.1))
    chip.adjustments[0] = 0.12
    chip.fill.solid()
    chip.fill.fore_color.rgb = TEAL_DIM if i in (0, 3) else BG_ALT
    chip.line.color.rgb = TEAL
    chip.line.width = Pt(1.25)
    chip.shadow.inherit = False
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = item
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK; r.font.name = FONT
    if i < len(py_items) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, py_bx + py_bw, py_dy + Inches(0.38), py_gap, Inches(0.34))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL
        arrow.line.fill.background(); arrow.shadow.inherit = False
    py_bx += py_bw + py_gap
add_bullets(s, Inches(0.7), Inches(5.15), Inches(11.7), Inches(1.6), [
    "Fourier Transforms: the FFT math behind every spectrum plot and "
    "waterfall shown tonight.",
    "Filters: the same low-pass/channel-select filtering every flowgraph "
    "tonight used, written out by hand.",
    "Digital Modulation: ASK, PSK, QAM, FSK -- how bits become RF, from "
    "first principles.",
    "RX & TX in Python: PlutoSDR, USRP, RTL-SDR, and HackRF, each driven "
    "directly from a Python script -- no GRC canvas required.",
], size=13, gap=4)
add_text(s, Inches(0.7), Inches(7.08), Inches(9), Inches(0.3),
          "pysdr.org -- \"A Guide to SDR and DSP using Python\"",
          size=10, color=MUTED, italic=True)
page_number(s, count(), TOTAL_SLIDES)

# ============================================================= SLIDE 24 ==
s = new_slide()
add_text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2),
          "Thank You", size=48, color=INK, bold=True)
add_text(s, Inches(0.9), Inches(4.0), Inches(11.0), Inches(0.6),
          "Questions? Let's get on the air.", size=20, color=BODY)
accent_bar(s, y=Inches(4.6), x=Inches(0.92), w=Inches(1.4))
add_text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4),
          "Mark Grennan  •  W5TSU  •  mark@grennan.com", size=14, color=MUTED)

prs.save("SDR_Demo.pptx")
print(f"Wrote SDR_Demo.pptx with {len(prs.slides._sldIdLst)} slides")
